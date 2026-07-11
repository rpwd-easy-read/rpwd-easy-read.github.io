"""Build RPwD_Act_2016_Easy_Read_Guide.pptx — Easy Read sample of the
Rights of Persons with Disabilities Act 2016 for persons with intellectual
and learning disabilities.

Each Section slide preserves the official Act title and section number so
the reader can cite the law when raising a complaint. Plain-language
captions and friendly comic-strip illustrations sit beneath the legal
heading, never in place of it.

Visual system:
- Vidyasagar palette (blue #00B0F0 + pink #F080B0) on a cream canvas
- A distinct chapter colour for each of the 17 chapters — for memorability,
  every section in a chapter shares its top band, badge and "use this when"
  colour, so a reader can recognise which chapter they're in at a glance
- Verdana throughout (Easy-Read standard sans-serif, ships with Office)
- Large body text, plenty of white space, one idea per slide
- PNG illustrations rendered by `_build/illustrations.py` (run that first)

Section titles and chapter assignments were extracted from the official
RPwD Act 2016 PDF (Ministry of Law and Justice text) — see
_build/section_titles.json. All 102 sections across 17 chapters included.

Run: python _build/illustrations.py    # generate PNGs first
     python _build/build_rpwd_easy_read.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ============================================================
# palette — Vidyasagar + per-chapter Easy-Read colours
# ============================================================

# Vidyasagar brand (sampled from the logo)
VS_BLUE       = RGBColor(0x00, 0xB0, 0xF0)
VS_BLUE_DARK  = RGBColor(0x00, 0x7C, 0xB0)
VS_PINK       = RGBColor(0xF0, 0x80, 0xB0)
VS_PINK_DARK  = RGBColor(0xC8, 0x5C, 0x90)

# Neutrals
CREAM         = RGBColor(0xFF, 0xFA, 0xF0)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
PAPER         = RGBColor(0xFF, 0xFF, 0xFF)
DARK          = RGBColor(0x1A, 0x1A, 0x1A)
MUTED         = RGBColor(0x55, 0x55, 0x55)
LIGHT_GREY    = RGBColor(0xE6, 0xE6, 0xE6)

# 17 chapter colours — distinct, all dark enough for white text
CHAPTER_COLOR = {
    "I":    RGBColor(0x00, 0xAC, 0xC1),  # cyan-blue
    "II":   RGBColor(0xC2, 0x18, 0x5B),  # pink
    "III":  RGBColor(0xF5, 0x7C, 0x00),  # orange
    "IV":   RGBColor(0x2E, 0x7D, 0x32),  # green
    "V":    RGBColor(0xC6, 0x28, 0x28),  # red
    "VI":   RGBColor(0x6A, 0x1B, 0x9A),  # purple
    "VII":  RGBColor(0x00, 0x89, 0x7B),  # teal
    "VIII": RGBColor(0x28, 0x35, 0x93),  # indigo
    "IX":   RGBColor(0x5D, 0x40, 0x37),  # brown
    "X":    RGBColor(0x00, 0x83, 0x8F),  # deep cyan
    "XI":   RGBColor(0x51, 0x2D, 0xA8),  # deep purple
    "XII":  RGBColor(0x68, 0x9F, 0x38),  # light green
    "XIII": RGBColor(0xBF, 0x36, 0x0C),  # deep orange
    "XIV":  RGBColor(0x45, 0x5A, 0x64),  # blue-grey
    "XV":   RGBColor(0x4E, 0x34, 0x2E),  # dark brown
    "XVI":  RGBColor(0xB7, 0x1C, 0x1C),  # dark red
    "XVII": RGBColor(0x42, 0x42, 0x42),  # dark grey
}


# ============================================================
# paths and constants
# ============================================================

PROJECT          = r"C:\Users\deepa\Projects\rpwd-easy-read"
OUTPUT           = os.path.join(PROJECT, "RPwD_Act_2016_Easy_Read_Guide.pptx")
LOGO             = os.path.join(PROJECT, "assets", "vidyasagar_logo.png")
ILLUSTRATION_DIR = os.path.join(PROJECT, "assets", "illustrations")

SW = Inches(13.333)
SH = Inches(7.5)
LABEL = "RPwD Act 2016  \u00b7  Easy Read Guide  \u00b7  All 102 Sections"

FONT_HEADING = "Verdana"
FONT_BODY    = "Verdana"


# ============================================================
# helpers
# ============================================================

def _style(run, *, name=FONT_BODY, size=18, color=DARK, bold=False, italic=False):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def set_bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, *,
             name=FONT_BODY, size=18, color=DARK, bold=False, italic=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    _style(r, name=name, size=size, color=color, bold=bold, italic=italic)
    return tb


def add_bullets(slide, left, top, width, height, items, *,
                size=18, color=DARK, name=FONT_BODY,
                bullet="\u2022", line_after=10, bold=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(line_after)
        indent = "    " * level
        prefix = f"{bullet}  " if level == 0 else "\u2013  "
        r = p.add_run()
        r.text = f"{indent}{prefix}{text}"
        _style(r, size=size, color=color, name=name, bold=bold)
    return tb


def add_rect(slide, left, top, width, height, fill, *, line=None, line_w=0.75):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    return s


def add_round_rect(slide, left, top, width, height, fill, *, line=None, line_w=0.75):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.adjustments[0] = 0.2
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    return s


def add_circle(slide, left, top, size, fill, *, line=None, line_w=0.75):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    return s


def footer(slide, num, total, color=VS_BLUE_DARK):
    bar_h = Inches(0.32)
    add_rect(slide, Inches(0), SH - bar_h, SW, bar_h, color)
    add_text(slide, Inches(0.4), SH - bar_h, Inches(10), bar_h,
             LABEL, size=10, color=WHITE, anchor=MSO_ANCHOR.MIDDLE,
             name=FONT_BODY)
    add_text(slide, SW - Inches(1.7), SH - bar_h, Inches(1.3), bar_h,
             f"Page {num} of {total}", size=10, color=WHITE,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
             name=FONT_BODY)


def new_blank(prs, bg=CREAM):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, bg)
    return slide


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def add_logo(slide, left, top, width, height, *, on_dark=False):
    """Vidyasagar logo with URL underneath."""
    if os.path.exists(LOGO) and os.path.getsize(LOGO) > 1024:
        ratio = 254 / 232
        logo_h = height - Inches(0.3)
        logo_w = int(logo_h * ratio)
        logo_left = left + (width - logo_w) // 2
        try:
            slide.shapes.add_picture(LOGO, logo_left, top, width=logo_w, height=logo_h)
            url_color = WHITE if on_dark else MUTED
            add_text(slide, left, top + logo_h, width, Inches(0.3),
                     "vidyasagar.co.in",
                     size=10, italic=True, color=url_color, align=PP_ALIGN.CENTER)
            return
        except Exception:
            pass
    color = WHITE if on_dark else MUTED
    add_text(slide, left, top, width, Inches(0.3),
             "Supported by VIDYASAGAR \u00b7 vidyasagar.co.in",
             size=11, italic=True, color=color, align=PP_ALIGN.CENTER)


def section_number_badge(slide, sec_text, color, x, y, size=Inches(1.0)):
    """White circle with section number, ringed in chapter color."""
    add_circle(slide, x, y, size, WHITE, line=color, line_w=4)
    add_text(slide, x, y, size, size, sec_text,
             name=FONT_HEADING, size=22, bold=True, color=color,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ============================================================
# slide builders — frontmatter
# ============================================================

def s_cover(prs, num, total):
    slide = new_blank(prs, bg=CREAM)

    # Two-tone top band: blue left, pink right
    add_rect(slide, Inches(0), Inches(0), SW // 2, Inches(0.55), VS_BLUE)
    add_rect(slide, SW // 2, Inches(0), SW // 2, Inches(0.55), VS_PINK)

    # Big title
    add_text(slide, Inches(0.7), Inches(1.1), Inches(12.0), Inches(1.4),
             "The RPwD Act 2016",
             name=FONT_HEADING, size=58, bold=True, color=VS_BLUE_DARK,
             anchor=MSO_ANCHOR.MIDDLE)

    # Subtitle in pink
    add_text(slide, Inches(0.7), Inches(2.55), Inches(12.0), Inches(0.8),
             "Easy Read Guide  \u00b7  All 102 Sections",
             name=FONT_HEADING, size=30, bold=True, italic=True, color=VS_PINK_DARK,
             anchor=MSO_ANCHOR.MIDDLE)

    # Decorative band (thin)
    add_rect(slide, Inches(0.7), Inches(3.55), Inches(12.0), Inches(0.06), VS_BLUE)

    # Tagline
    add_text(slide, Inches(0.7), Inches(3.85), Inches(12.0), Inches(0.7),
             "Your rights, in plain words \u2014 with the law still attached.",
             name=FONT_BODY, size=24, color=DARK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Audience
    add_text(slide, Inches(0.7), Inches(4.6), Inches(12.0), Inches(0.5),
             "A complete reference for persons with intellectual and learning disabilities.",
             name=FONT_BODY, size=17, italic=True, color=MUTED,
             align=PP_ALIGN.CENTER)

    # Verification line
    add_text(slide, Inches(0.7), Inches(5.15), Inches(12.0), Inches(0.4),
             "Section titles verified against the official RPwD Act 2016 (Ministry of Law and Justice).",
             name=FONT_BODY, size=12, italic=True, color=MUTED,
             align=PP_ALIGN.CENTER)

    # Vidyasagar attribution — bottom right, subtle
    add_logo(slide, Inches(10.9), Inches(5.95), Inches(2.0), Inches(1.05))

    # Bottom accent
    add_rect(slide, Inches(0), Inches(7.05), SW, Inches(0.15), VS_PINK)

    footer(slide, num, total, color=VS_BLUE_DARK)
    notes(slide, "Cover. Title: The RPwD Act 2016 — Easy Read Guide. All 102 sections. Section titles verified against the official Ministry of Law and Justice text. Supported by Vidyasagar.")


def s_about(prs, num, total):
    slide = new_blank(prs, bg=CREAM)
    add_rect(slide, Inches(0), Inches(0), SW, Inches(0.55), VS_BLUE)
    add_text(slide, Inches(0.5), Inches(0.05), Inches(12), Inches(0.5),
             "About this guide",
             name=FONT_HEADING, size=24, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)

    add_bullets(slide, Inches(0.7), Inches(1.1), Inches(11.9), Inches(4.6), [
        "This guide explains all 102 sections of India's Rights of Persons with Disabilities Act 2016.",
        "It is written in plain words for people who find legal language hard.",
        "Every page keeps the official Section number and title from the Act. The plain words sit below.",
        "Some sections give you rights you can claim. Others explain how the system works \u2014 those say \u201CFor reference\u201D in the Use this when box.",
        "Each chapter has its own colour, so you can tell at a glance which chapter you are in.",
        "Section titles in this guide were verified against the official Act text from the Ministry of Law and Justice.",
        "This guide is not legal advice. For court matters, please ask a lawyer or your State Commissioner.",
    ], size=18, line_after=12)

    add_round_rect(slide, Inches(0.5), Inches(5.95), Inches(12.3), Inches(1.0), VS_PINK)
    add_text(slide, Inches(0.85), Inches(6.0), Inches(11.6), Inches(0.4),
             "WHY THE LEGAL HEADINGS STAY",
             name=FONT_HEADING, size=12, bold=True, color=WHITE)
    add_text(slide, Inches(0.85), Inches(6.35), Inches(11.6), Inches(0.6),
             "When you raise a complaint, citing the Section number gives your case real power. We never strip it out.",
             name=FONT_BODY, size=15, color=WHITE)

    footer(slide, num, total)
    notes(slide, "About this guide. Plain explanation of all 102 sections of the RPwD Act 2016. Section titles verified against the official Act text. Not legal advice.")


def s_how_to_use(prs, num, total):
    slide = new_blank(prs, bg=CREAM)
    add_rect(slide, Inches(0), Inches(0), SW, Inches(0.55), VS_BLUE)
    add_text(slide, Inches(0.5), Inches(0.05), Inches(12), Inches(0.5),
             "How to use this guide",
             name=FONT_HEADING, size=24, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)

    add_text(slide, Inches(0.7), Inches(1.0), Inches(11.9), Inches(0.5),
             "Every Section page in this guide has the same five parts:",
             name=FONT_BODY, size=18, italic=True, color=MUTED)

    items = [
        ("1", "Chapter and Section number", "Tells you exactly where in the Act you are.", VS_BLUE),
        ("2", "Official Section title",     "The exact name from the Act. Use this in complaints.", VS_PINK),
        ("3", "Picture",                    "A simple drawing to help you remember the idea.", VS_BLUE),
        ("4", "Plain words",                "The same idea in easy English.", VS_PINK),
        ("5", "Use this when…",             "Real-life examples of when this section helps you.", VS_BLUE),
    ]

    card_y = Inches(1.7)
    card_h = Inches(0.85)
    gap = Inches(0.12)
    for i, (n, head_t, body_t, color) in enumerate(items):
        y = card_y + i * (card_h + gap)
        add_round_rect(slide, Inches(0.5), y, Inches(0.95), card_h, color)
        add_text(slide, Inches(0.5), y, Inches(0.95), card_h, n,
                 name=FONT_HEADING, size=28, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_round_rect(slide, Inches(1.65), y, Inches(11.2), card_h, WHITE,
                       line=color, line_w=2)
        add_text(slide, Inches(1.85), y + Inches(0.08), Inches(10.9), Inches(0.36),
                 head_t, name=FONT_HEADING, size=15, bold=True, color=color,
                 anchor=MSO_ANCHOR.TOP)
        add_text(slide, Inches(1.85), y + Inches(0.42), Inches(10.9), Inches(0.42),
                 body_t, name=FONT_BODY, size=13, color=DARK,
                 anchor=MSO_ANCHOR.TOP)

    footer(slide, num, total)
    notes(slide, "How to use this guide. The five-part anatomy of every Section page. Legal anchors stay so the reader can cite them in complaints.")


def s_map(prs, num, total):
    slide = new_blank(prs, bg=CREAM)
    add_rect(slide, Inches(0), Inches(0), SW, Inches(0.55), VS_BLUE)
    add_text(slide, Inches(0.5), Inches(0.05), Inches(12), Inches(0.5),
             "Map of your rights",
             name=FONT_HEADING, size=24, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)

    add_text(slide, Inches(0.7), Inches(0.95), Inches(11.9), Inches(0.4),
             "All 17 chapters and 102 sections \u2014 each chapter has its own colour:",
             name=FONT_BODY, size=14, italic=True, color=MUTED)

    rows = [
        ("CH I",    "The basics",                              "1, 2",            "I"),
        ("CH II",   "Your rights",                              "3 \u2013 15",     "II"),
        ("CH III",  "Education",                                "16 \u2013 18",    "III"),
        ("CH IV",   "Skill and work",                           "19 \u2013 23",    "IV"),
        ("CH V",    "Health and support",                       "24 \u2013 30",    "V"),
        ("CH VI",   "Benchmark disabilities",                   "31 \u2013 37",    "VI"),
        ("CH VII",  "High support needs",                       "38",              "VII"),
        ("CH VIII", "Government duties",                        "39 \u2013 48",    "VIII"),
        ("CH IX",   "Disability institutions",                  "49 \u2013 55",    "IX"),
        ("CH X",    "Certification",                            "56 \u2013 59",    "X"),
        ("CH XI",   "Advisory boards",                          "60 \u2013 73",    "XI"),
        ("CH XII",  "Commissioners",                            "74 \u2013 83",    "XII"),
        ("CH XIII", "Special Court",                            "84, 85",          "XIII"),
        ("CH XIV",  "National Fund",                            "86, 87",          "XIV"),
        ("CH XV",   "State Fund",                               "88",              "XV"),
        ("CH XVI",  "Penalties",                                "89 \u2013 95",    "XVI"),
        ("CH XVII", "Miscellaneous",                            "96 \u2013 102",   "XVII"),
    ]

    chip_w = Inches(1.05)
    name_w = Inches(3.5)
    secs_w = Inches(1.45)
    row_h = Inches(0.42)
    gap = Inches(0.06)

    left_x = Inches(0.5)
    right_x = Inches(6.85)
    start_y = Inches(1.45)

    for i, (ch, n, secs, ch_key) in enumerate(rows):
        col_x = left_x if i < 9 else right_x
        local_i = i if i < 9 else i - 9
        y = start_y + local_i * (row_h + gap)
        col = CHAPTER_COLOR[ch_key]
        add_round_rect(slide, col_x, y, chip_w, row_h, col)
        add_text(slide, col_x, y, chip_w, row_h, ch,
                 name=FONT_HEADING, size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, col_x + chip_w + Inches(0.12), y, name_w, row_h, n,
                 name=FONT_HEADING, size=14, bold=True, color=DARK,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, col_x + chip_w + Inches(0.12) + name_w, y, secs_w, row_h, secs,
                 name=FONT_BODY, size=12, italic=True, color=col,
                 bold=True, anchor=MSO_ANCHOR.MIDDLE)

    add_text(slide, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
             "Plus: How to file a complaint  \u00b7  Who can help you  \u00b7  Your rights in one page",
             name=FONT_BODY, size=14, italic=True, color=MUTED,
             align=PP_ALIGN.CENTER)

    footer(slide, num, total)
    notes(slide, "Map of your rights. All 17 chapters and 102 sections of the RPwD Act 2016, with the section number range and the chapter colour for each.")


# ============================================================
# slide builders — chapter dividers and section pages
# ============================================================

def s_chapter_divider(prs, num, total, *, ch, kicker, title, subtitle):
    color = CHAPTER_COLOR[ch]
    slide = new_blank(prs, bg=color)

    # decorative top stripe
    add_rect(slide, Inches(0), Inches(0), SW, Inches(0.4), WHITE)
    add_rect(slide, Inches(0), Inches(0.4), SW, Inches(0.08), VS_PINK)

    add_text(slide, Inches(0.8), Inches(2.0), Inches(11.7), Inches(0.6),
             kicker, name=FONT_BODY, size=22, color=WHITE, italic=True,
             align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.8), Inches(2.8), Inches(11.7), Inches(1.7),
             title, name=FONT_HEADING, size=44, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(0.8), Inches(4.85), Inches(11.7), Inches(0.8),
             subtitle, name=FONT_BODY, size=20, italic=True, color=WHITE,
             align=PP_ALIGN.CENTER)

    # decorative bottom stripe
    add_rect(slide, Inches(0), SH - Inches(0.55), SW, Inches(0.08), VS_PINK)
    add_rect(slide, Inches(0), SH - Inches(0.47), SW, Inches(0.15), WHITE)

    footer(slide, num, total, color=color)
    notes(slide, f"{kicker}: {title}. {subtitle}")


def s_section(prs, num, total, *,
              chapter, section_title, plain_text, use_when, illustration,
              speaker_notes=None, title_size=26):
    color = CHAPTER_COLOR[chapter]
    chapter_label = CHAPTER_LABEL[chapter]
    slide = new_blank(prs, bg=CREAM)

    # Top chapter band
    band_h = Inches(0.65)
    add_rect(slide, Inches(0), Inches(0), SW, band_h, color)
    add_text(slide, Inches(0.5), Inches(0), Inches(10), band_h,
             chapter_label,
             name=FONT_HEADING, size=15, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)

    # Section number badge — overlapping the band
    sec_num = section_title.split()[1]  # "Section 3 — ..." -> "3"
    badge_size = Inches(1.05)
    section_number_badge(slide, sec_num, color,
                         x=SW - badge_size - Inches(0.55),
                         y=Inches(0.1),
                         size=badge_size)

    # Illustration on the left (square)
    img_path = os.path.join(ILLUSTRATION_DIR, f"{illustration}.png")
    img_size = Inches(4.0)
    img_x = Inches(0.5)
    img_y = Inches(1.15)
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, img_x, img_y,
                                 width=img_size, height=img_size)
    else:
        # placeholder
        add_round_rect(slide, img_x, img_y, img_size, img_size, LIGHT_GREY)
        add_text(slide, img_x, img_y, img_size, img_size,
                 f"[{illustration}]",
                 size=14, italic=True, color=MUTED,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Heading on the right
    head_x = Inches(4.85)
    head_w = Inches(8.0)
    add_text(slide, head_x, Inches(1.15), head_w, Inches(1.5),
             section_title,
             name=FONT_HEADING, size=title_size, bold=True, color=color)

    # Plain language caption
    add_text(slide, head_x, Inches(2.75), head_w, Inches(2.55),
             plain_text,
             name=FONT_BODY, size=20, color=DARK)

    # USE THIS WHEN box (full width)
    box_top = Inches(5.5)
    box_h = Inches(1.5)
    add_round_rect(slide, Inches(0.5), box_top, Inches(12.3), box_h, color)
    add_text(slide, Inches(0.85), box_top + Inches(0.12), Inches(11.6), Inches(0.4),
             "\u2605  USE THIS WHEN",
             name=FONT_HEADING, size=14, bold=True, color=WHITE)
    add_text(slide, Inches(0.85), box_top + Inches(0.55), Inches(11.6), Inches(0.95),
             use_when,
             name=FONT_BODY, size=17, color=WHITE)

    footer(slide, num, total, color=color)
    sn = speaker_notes or f"{section_title}. Plain meaning: {plain_text} Use when: {use_when}"
    notes(slide, sn)
    return slide


# ============================================================
# slide builders — backmatter
# ============================================================

def s_complain(prs, num, total):
    slide = new_blank(prs, bg=CREAM)
    add_rect(slide, Inches(0), Inches(0), SW, Inches(0.65), VS_BLUE)
    add_text(slide, Inches(0.5), Inches(0), Inches(12), Inches(0.65),
             "How to file a complaint",
             name=FONT_HEADING, size=26, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)

    add_text(slide, Inches(0.7), Inches(1.0), Inches(11.9), Inches(0.4),
             "Five steps to raise a complaint under the RPwD Act:",
             name=FONT_BODY, size=17, italic=True, color=MUTED)

    steps = [
        ("1", "Write down what happened.",
         "Date, place, names if you know them. Two or three lines is fine.",
         VS_BLUE),
        ("2", "Find the right Section number.",
         "Use this guide to pick the Section that fits your problem. Write it on your complaint.",
         VS_PINK),
        ("3", "Send your complaint.",
         "Send it to the State Commissioner for Persons with Disabilities in your state. By email or by post.",
         VS_BLUE),
        ("4", "Keep proof.",
         "Save a copy. Get a receipt, an email reply, or a stamp showing they got your complaint.",
         VS_PINK),
        ("5", "Wait 60 days.",
         "If nothing happens in 60 days, you can take it to the Chief Commissioner or to court.",
         VS_BLUE),
    ]

    y = Inches(1.55)
    h = Inches(1.0)
    for n, head_t, body_t, col in steps:
        # number circle
        circ_size = Inches(0.85)
        circ_x = Inches(0.7)
        circ_y = y + (h - circ_size) // 2
        add_circle(slide, circ_x, circ_y, circ_size, col)
        add_text(slide, circ_x, circ_y, circ_size, circ_size, n,
                 name=FONT_HEADING, size=28, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # rounded card
        add_round_rect(slide, Inches(1.85), y, Inches(11.0), h, WHITE,
                       line=col, line_w=2)
        add_text(slide, Inches(2.1), y + Inches(0.1), Inches(10.7), Inches(0.45),
                 head_t, name=FONT_HEADING, size=18, bold=True, color=col)
        add_text(slide, Inches(2.1), y + Inches(0.55), Inches(10.7), Inches(0.42),
                 body_t, name=FONT_BODY, size=14, color=DARK)
        y += h + Inches(0.08)

    footer(slide, num, total)
    notes(slide, "How to file a complaint: 5 steps from writing it down to escalating to the Chief Commissioner if no response in 60 days.")


def s_help(prs, num, total):
    slide = new_blank(prs, bg=CREAM)
    add_rect(slide, Inches(0), Inches(0), SW, Inches(0.65), VS_BLUE)
    add_text(slide, Inches(0.5), Inches(0), Inches(12), Inches(0.65),
             "Who can help you",
             name=FONT_HEADING, size=26, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)

    add_text(slide, Inches(0.7), Inches(1.0), Inches(11.9), Inches(0.4),
             "If your rights are not respected, these offices can help:",
             name=FONT_BODY, size=17, italic=True, color=MUTED)

    cards = [
        ("State Commissioner for Persons with Disabilities",
         "Every state in India has one. They handle complaints inside the state. Look up your state's office or ask the state social welfare department.",
         "First place to go for most complaints.",
         VS_BLUE),
        ("Chief Commissioner for Persons with Disabilities",
         "A national office for cases that go beyond one state, or when the State Commissioner has not helped within 60 days.",
         "Office of the Chief Commissioner, New Delhi.",
         VS_PINK),
        ("National Human Rights Commission",
         "For very serious cases \u2014 cruelty, abuse, denial of basic rights.",
         "Especially under Sections 6, 7, and 92.",
         VS_BLUE),
        ("Disabled People's Organisations",
         "Local groups led by persons with disabilities. They can help you write your complaint and stand with you.",
         "Vidyasagar in Chennai is one such organisation.",
         VS_PINK),
    ]

    card_w = Inches(6.0)
    card_h = Inches(2.45)
    gap_x = Inches(0.3)
    gap_y = Inches(0.2)
    start_x = Inches(0.5)
    start_y = Inches(1.55)
    for i, (head_t, body_t, foot_t, col) in enumerate(cards):
        c = i % 2
        r = i // 2
        x = start_x + c * (card_w + gap_x)
        y = start_y + r * (card_h + gap_y)
        # outer card
        add_round_rect(slide, x, y, card_w, card_h, WHITE, line=col, line_w=2.5)
        # head bar
        add_round_rect(slide, x, y, card_w, Inches(0.55), col)
        add_text(slide, x + Inches(0.15), y, card_w - Inches(0.3), Inches(0.55),
                 head_t, name=FONT_HEADING, size=14, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
        # body
        add_text(slide, x + Inches(0.15), y + Inches(0.7),
                 card_w - Inches(0.3), Inches(1.25),
                 body_t, name=FONT_BODY, size=13, color=DARK)
        # foot
        add_text(slide, x + Inches(0.15), y + card_h - Inches(0.5),
                 card_w - Inches(0.3), Inches(0.4),
                 foot_t, name=FONT_BODY, size=12, italic=True, color=col)

    footer(slide, num, total)
    notes(slide, "Who can help you: State Commissioner, Chief Commissioner, NHRC, and Disabled People's Organisations. Four main avenues for raising a disability rights complaint in India.")


def s_summary(prs, num, total):
    slide = new_blank(prs, bg=CREAM)
    add_rect(slide, Inches(0), Inches(0), SW, Inches(0.65), VS_BLUE)
    add_text(slide, Inches(0.5), Inches(0), Inches(12), Inches(0.65),
             "Remember: your rights in one page",
             name=FONT_HEADING, size=24, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)

    points = [
        ("You are equal.",                            "No one can treat you badly.",              "Section 3",     VS_BLUE),
        ("You can live with your family.",            "Not in an institution against your will.", "Section 5",     VS_PINK),
        ("No one can hurt or use you.",               "You can call the police.",                  "Sections 6, 7", VS_BLUE),
        ("You can go to school.",                     "Schools must make space for you.",         "Sections 16, 17", VS_PINK),
        ("You can work.",                             "4 in 100 government jobs are reserved.",   "Sections 20, 34", VS_BLUE),
        ("You can have healthcare and a pension.",    "Free, close to home.",                     "Sections 24, 25", VS_PINK),
        ("Buildings and information must work for you.", "Ramps. Lifts. Easy formats.",          "Sections 40, 42", VS_BLUE),
        ("If anyone breaks these rules,",             "they can be fined or jailed.",             "Sections 89, 92", VS_PINK),
    ]

    y = Inches(1.15)
    h = Inches(0.6)
    for head_t, body_t, sec, col in points:
        add_circle(slide, Inches(0.7), y + Inches(0.12), Inches(0.36), col)
        add_text(slide, Inches(1.25), y, Inches(6.2), h, head_t,
                 name=FONT_HEADING, size=17, bold=True, color=DARK,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(7.45), y, Inches(4.2), h, body_t,
                 name=FONT_BODY, size=15, color=DARK,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(11.7), y, Inches(1.55), h, sec,
                 name=FONT_BODY, size=12, italic=True, bold=True, color=col,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        y += h + Inches(0.05)

    add_round_rect(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.7), VS_PINK)
    add_text(slide, Inches(0.85), Inches(6.4), Inches(11.6), Inches(0.7),
             "Always write the Section number when you raise a complaint. That is what makes it powerful.",
             name=FONT_HEADING, size=15, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    footer(slide, num, total)
    notes(slide, "Your rights in one page. Eight-line summary of the most-cited sections. Closing reminder: always write the Section number when raising a complaint.")


def s_back(prs, num, total):
    slide = new_blank(prs, bg=VS_BLUE)

    # Pink decorative band
    add_rect(slide, Inches(0), Inches(0.95), SW, Inches(0.12), VS_PINK)

    add_text(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.0),
             "\u201CNothing about us without us\u201D",
             name=FONT_HEADING, size=38, italic=True, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_text(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(0.5),
             "Capacity, not dependency.",
             name=FONT_BODY, size=22, italic=True, color=WHITE,
             align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(5.5), Inches(3.3), Inches(2.3), Inches(0.06), WHITE)

    add_text(slide, Inches(0.8), Inches(3.6), Inches(11.7), Inches(0.6),
             "RPwD Act 2016 \u2014 Easy Read Guide",
             name=FONT_HEADING, size=26, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)

    add_text(slide, Inches(0.8), Inches(4.2), Inches(11.7), Inches(0.5),
             "All 102 Sections  \u00b7  April 2026",
             name=FONT_BODY, size=16, italic=True, color=WHITE,
             align=PP_ALIGN.CENTER)

    add_text(slide, Inches(0.8), Inches(4.7), Inches(11.7), Inches(0.5),
             "Prepared by the Disability Law Unit of Vidya Sagar  \u00b7  In support of Ummul Khair Shamim",
             name=FONT_BODY, size=14, color=WHITE,
             align=PP_ALIGN.CENTER)

    add_logo(slide, Inches(5.4), Inches(5.3), Inches(2.5), Inches(1.4), on_dark=True)

    add_text(slide, Inches(0.8), Inches(6.85), Inches(11.7), Inches(0.4),
             "Section titles verified against the official Ministry of Law and Justice text. Not legal advice.",
             name=FONT_BODY, size=12, italic=True, color=WHITE,
             align=PP_ALIGN.CENTER)

    # Bottom pink decorative band
    add_rect(slide, Inches(0), SH - Inches(0.47), SW, Inches(0.15), VS_PINK)

    footer(slide, num, total)
    notes(slide, "Back cover. Nothing about us without us. Capacity not dependency. All 102 sections. April 2026. Prepared by the Disability Law Unit of Vidya Sagar, in support of Ummul Khair Shamim.")


# ============================================================
# data — chapter dividers, chapter labels, sections
# ============================================================

CHAPTER_DIVIDERS = {
    "I":    dict(ch="I",    kicker="CHAPTER I",    title="Preliminary",
                 subtitle="The basics: who is covered and what the words mean"),
    "II":   dict(ch="II",   kicker="CHAPTER II",   title="Rights and Entitlements",
                 subtitle="What you can claim, and what is wrong"),
    "III":  dict(ch="III",  kicker="CHAPTER III",  title="Education",
                 subtitle="Schools, colleges, and learning"),
    "IV":   dict(ch="IV",   kicker="CHAPTER IV",   title="Skill Development and Employment",
                 subtitle="Training, jobs, and fair treatment at work"),
    "V":    dict(ch="V",    kicker="CHAPTER V",    title="Social Security, Health, Rehabilitation and Recreation",
                 subtitle="Help to live with dignity, stay well, and take part"),
    "VI":   dict(ch="VI",   kicker="CHAPTER VI",   title="Special Provisions for Persons with Benchmark Disabilities",
                 subtitle="Extra rights when your disability is 40% or more"),
    "VII":  dict(ch="VII",  kicker="CHAPTER VII",  title="Persons with Disabilities with High Support Needs",
                 subtitle="Extra support for those who need it most"),
    "VIII": dict(ch="VIII", kicker="CHAPTER VIII", title="Duties and Responsibilities of Appropriate Governments",
                 subtitle="Awareness, accessibility, and what governments must build"),
    "IX":   dict(ch="IX",   kicker="CHAPTER IX",   title="Registration of Disability Institutions",
                 subtitle="Rules for organisations that work with persons with disabilities"),
    "X":    dict(ch="X",    kicker="CHAPTER X",    title="Certification of Specified Disabilities",
                 subtitle="How to assess and certify a disability"),
    "XI":   dict(ch="XI",   kicker="CHAPTER XI",   title="Advisory Boards on Disability",
                 subtitle="Central, state, and district advisory bodies"),
    "XII":  dict(ch="XII",  kicker="CHAPTER XII",  title="Chief Commissioner and State Commissioners",
                 subtitle="The offices that handle your complaints"),
    "XIII": dict(ch="XIII", kicker="CHAPTER XIII", title="Special Court",
                 subtitle="Where serious cases under this Act are heard"),
    "XIV":  dict(ch="XIV",  kicker="CHAPTER XIV",  title="National Fund for Persons with Disabilities",
                 subtitle="The national pot of money that pays for support"),
    "XV":   dict(ch="XV",   kicker="CHAPTER XV",   title="State Fund for Persons with Disabilities",
                 subtitle="Each state's own disability fund"),
    "XVI":  dict(ch="XVI",  kicker="CHAPTER XVI",  title="Offences and Penalties",
                 subtitle="If someone breaks these rules"),
    "XVII": dict(ch="XVII", kicker="CHAPTER XVII", title="Miscellaneous",
                 subtitle="The closing rules of the Act"),
}


CHAPTER_LABEL = {
    "I":    "CHAPTER I  \u00b7  THE BASICS",
    "II":   "CHAPTER II  \u00b7  YOUR RIGHTS",
    "III":  "CHAPTER III  \u00b7  EDUCATION",
    "IV":   "CHAPTER IV  \u00b7  SKILL & WORK",
    "V":    "CHAPTER V  \u00b7  HEALTH AND SUPPORT",
    "VI":   "CHAPTER VI  \u00b7  BENCHMARK DISABILITIES",
    "VII":  "CHAPTER VII  \u00b7  HIGH SUPPORT NEEDS",
    "VIII": "CHAPTER VIII  \u00b7  GOVERNMENT DUTIES",
    "IX":   "CHAPTER IX  \u00b7  DISABILITY INSTITUTIONS",
    "X":    "CHAPTER X  \u00b7  CERTIFICATION",
    "XI":   "CHAPTER XI  \u00b7  ADVISORY BOARDS",
    "XII":  "CHAPTER XII  \u00b7  COMMISSIONERS",
    "XIII": "CHAPTER XIII  \u00b7  SPECIAL COURT",
    "XIV":  "CHAPTER XIV  \u00b7  NATIONAL FUND",
    "XV":   "CHAPTER XV  \u00b7  STATE FUND",
    "XVI":  "CHAPTER XVI  \u00b7  PENALTIES",
    "XVII": "CHAPTER XVII  \u00b7  MISCELLANEOUS",
}


def S(chapter, title, plain, use_when, il, *, ts=26):
    d = dict(chapter=chapter, section_title=title,
             plain_text=plain, use_when=use_when, illustration=il)
    if ts != 26:
        d["title_size"] = ts
    return d


SECTIONS = [
    # ---------- Chapter I ----------
    S("I", "Section 1 \u2014 Short title and commencement",
      "This section gives the Act its name and tells when it started. The Act began on 19 April 2017 and applies to all of India.",
      "(For reference. This section is the Act's name plate.)",
      "book_simple"),
    S("I", "Section 2 \u2014 Definitions",
      "The Act lists 21 disabilities. If you have one of them, you are covered. If your disability is 40% or more, you are called a person with benchmark disability — you get extra rights.",
      "You are not sure if the Act covers you, or whether you have a benchmark disability. Look at the list of 21 disabilities and the 40% rule.",
      "book_simple"),

    # ---------- Chapter II ----------
    S("II", "Section 3 \u2014 Equality and non-discrimination",
      "You have the right to be treated equally. No one can treat you badly because you have a disability.",
      "Someone refused you a service, a job, or a place because of your disability. Cite Section 3 in your complaint.",
      "equality"),
    S("II", "Section 4 \u2014 Women and children with disabilities",
      "Women and children with disabilities have the same rights as everyone else. The government must protect them from extra harm.",
      "A girl or woman is denied school, work, or safety because of her disability. Or a child is being kept out of school.",
      "women_children"),
    S("II", "Section 5 \u2014 Community life",
      "You have the right to live in your own home and your own community. No one can force you into an institution against your will.",
      "Someone is forcing a person with a disability to live in a hostel or institution they do not want to live in.",
      "house_family"),
    S("II", "Section 6 \u2014 Protection from cruelty and inhuman treatment",
      "No one can hurt you or treat you in a cruel way. The government must keep you safe.",
      "A person with a disability is being treated cruelly at home, in school, in a hospital, or in care.",
      "shield"),
    S("II", "Section 7 \u2014 Protection from abuse, violence and exploitation",
      "No one can hit you, hurt you, or use you for their gain. You can tell the police. They must help. You can also ask a magistrate for protection.",
      "Someone is hurting, abusing, or exploiting a person with a disability.",
      "shield"),
    S("II", "Section 8 \u2014 Protection and safety",
      "The government must keep persons with disabilities safe in disasters, fires, floods, armed conflict, and other emergencies.",
      "A relief camp, shelter, or evacuation plan does not include or is not safe for persons with disabilities.",
      "disaster_safe"),
    S("II", "Section 9 \u2014 Home and family",
      "Persons with disabilities have the right to a home and a family. No child can be taken from a parent only because the parent has a disability.",
      "A child welfare worker tries to take your child away because of your disability.",
      "house_family"),
    S("II", "Section 10 \u2014 Reproductive rights",
      "Persons with disabilities have the right to decide about marriage and children. No one can force you to have an operation that stops you from having children.",
      "A doctor, family member, or institution wants to sterilise a person with a disability without their full free consent.",
      "heart"),
    S("II", "Section 11 \u2014 Accessibility in voting",
      "You have the right to vote in private. Polling stations must have ramps, easy ballots, and help if you need it.",
      "A polling station has no ramp, no Braille ballot, or no help for voters with disabilities.",
      "vote"),
    S("II", "Section 12 \u2014 Access to justice",
      "You have the right to go to court. Courts must let you take part. They must give you a sign language interpreter or a support person if you need one.",
      "A court, police station, or tribunal does not let you take part because of your disability.",
      "justice_scales"),
    S("II", "Section 13 \u2014 Legal capacity",
      "You can make your own decisions. You can own things. You can run a bank account. No one can take that away because of your disability.",
      "Someone says you cannot sign a document, open a bank account, or own property because of your disability.",
      "pencil_signing"),
    S("II", "Section 14 \u2014 Provision for guardianship",
      "Guardianship can only be a last step. It must be limited and end when you can decide for yourself again. You always have the right to support to make your own choices.",
      "Someone wants to put a person with a disability under full guardianship without trying support first.",
      "friends_support"),
    S("II", "Section 15 \u2014 Designation of authorities to support",
      "The government must name offices that help persons with disabilities make their own decisions.",
      "(For reference. This section sets up the support system for decision-making.)",
      "document"),

    # ---------- Chapter III ----------
    S("III", "Section 16 \u2014 Duty of educational institutions",
      "All schools must welcome children with disabilities. Schools must give the help each child needs to learn.",
      "A school refused to take a child because of disability, or refuses to give support.",
      "school"),
    S("III", "Section 17 \u2014 Specific measures to promote and facilitate inclusive education",
      "The government must train teachers, make books in accessible formats, and make schools easy to enter and use.",
      "A school does not have ramps, accessible books, sign language support, or trained teachers.",
      "book_open", ts=20),
    S("III", "Section 18 \u2014 Adult education",
      "The government must run adult education programmes that include adults with disabilities.",
      "An adult literacy or skills programme has no place for persons with disabilities.",
      "book_learning"),

    # ---------- Chapter IV ----------
    S("IV", "Section 19 \u2014 Vocational training and self-employment",
      "The government must run job training programmes and help persons with disabilities start their own work.",
      "A skill training programme refused you, or did not have an accessible course.",
      "briefcase"),
    S("IV", "Section 20 \u2014 Non-discrimination in employment",
      "A government job cannot say no to you because you have a disability. You must be treated fairly at work.",
      "You did not get a job, a promotion, or a fair review because of your disability.",
      "briefcase"),
    S("IV", "Section 21 \u2014 Equal opportunity policy",
      "Every workplace with 20 or more workers must write down how they will treat workers with disabilities fairly. They must show this paper if you ask.",
      "You want to know your workplace's rules for disability support. Ask to see the equal opportunity policy.",
      "document"),
    S("IV", "Section 22 \u2014 Maintenance of records",
      "Every workplace must keep records of how many workers with disabilities they have, and what support each one needs.",
      "You ask your workplace for these records and they refuse, or they say no records are kept.",
      "clipboard_check"),
    S("IV", "Section 23 \u2014 Appointment of Grievance Redressal Officer",
      "Every government workplace must name an officer to handle complaints from workers with disabilities.",
      "Your workplace has no Grievance Redressal Officer, or the officer is not doing the job.",
      "committee_people"),

    # ---------- Chapter V ----------
    S("V", "Section 24 \u2014 Social security",
      "The government must help you with food, clothes, a home, and other support if you need it to live.",
      "You need a disability pension, ration card, or housing scheme but were turned away.",
      "money_jar"),
    S("V", "Section 25 \u2014 Healthcare",
      "You should get free healthcare close to where you live. Healthcare must work for people with disabilities too.",
      "A government hospital or clinic refused you, charged you, or had no accessible service.",
      "medical_cross"),
    S("V", "Section 26 \u2014 Insurance schemes",
      "The government must run insurance schemes for persons with disabilities.",
      "An insurance company refused you because of your disability, or you cannot find any disability insurance.",
      "document"),
    S("V", "Section 27 \u2014 Rehabilitation",
      "The government must run rehabilitation programmes covering health, mind, social skills, and work.",
      "A government rehab centre is not available, is too far, or refused to take you.",
      "friends_support"),
    S("V", "Section 28 \u2014 Research and development",
      "The government must support research on disability, including new tools and treatments.",
      "(For reference. This section funds research on assistive technology and disability services.)",
      "book_open"),
    S("V", "Section 29 \u2014 Culture and recreation",
      "Persons with disabilities have the right to take part in art, music, films, libraries, and leisure.",
      "A museum, cinema, library, or cultural event refused you or had no accessible entry.",
      "heart"),
    S("V", "Section 30 \u2014 Sporting activities",
      "The government must support sports for persons with disabilities.",
      "A sports facility, training programme, or competition refused to include persons with disabilities.",
      "certificate_star"),

    # ---------- Chapter VI ----------
    S("VI", "Section 31 \u2014 Free education for children with benchmark disabilities",
      "Children with benchmark disabilities aged 6 to 18 have the right to free education in any government or government-aided school.",
      "A government school asked a child with benchmark disability to pay fees, or refused to admit them.",
      "school", ts=22),
    S("VI", "Section 32 \u2014 Reservation in higher educational institutions",
      "Government colleges and universities must keep 5 seats out of every 100 for students with benchmark disabilities.",
      "A college denied you a reserved seat or did not have a 5% disability quota in admissions.",
      "school"),
    S("VI", "Section 33 \u2014 Identification of posts for reservation",
      "The government must list out which jobs are open to persons with benchmark disabilities under the reservation.",
      "(For reference. This section is the basis for the 4% reservation in Section 34.)",
      "document"),
    S("VI", "Section 34 \u2014 Reservation",
      "The government must keep 4 jobs out of every 100 for persons with benchmark disabilities. The 4% is split across categories of disability.",
      "A government office did not follow the 4% job quota, or did not let you apply under the disability category.",
      "chair_reserved"),
    S("VI", "Section 35 \u2014 Incentives to employers in private sector",
      "The government can give private companies incentives like grants and tax breaks if they hire workers with disabilities.",
      "(For reference. This section encourages private sector hiring of persons with disabilities.)",
      "money_jar"),
    S("VI", "Section 36 \u2014 Special employment exchange",
      "The government must set up special job centres that help persons with disabilities find work.",
      "You cannot find a special employment exchange in your district, or the one you found is not helping you.",
      "briefcase"),
    S("VI", "Section 37 \u2014 Special schemes and development programmes",
      "The government must run schemes that help persons with benchmark disabilities live, study, and work.",
      "(For reference. This section is the basis for many disability welfare schemes.)",
      "document"),

    # ---------- Chapter VII ----------
    S("VII", "Section 38 \u2014 Special provisions for persons with disabilities with high support",
      "Persons with disabilities who need a lot of support get extra help. They have the right to ask for support, and the help must be reviewed often.",
      "A person needs daily personal support but is not getting it from any government scheme.",
      "high_support_hands", ts=20),

    # ---------- Chapter VIII ----------
    S("VIII", "Section 39 \u2014 Awareness campaigns",
      "The government must run programmes to teach people about disability rights and accessibility.",
      "(For reference. This section requires public awareness of the Act.)",
      "megaphone"),
    S("VIII", "Section 40 \u2014 Accessibility",
      "Buildings, places, transport, and information must be made easy to use for everyone, including persons with disabilities.",
      "A public building has no ramp, no lift, no accessible toilet, or no clear signs.",
      "building_ramp"),
    S("VIII", "Section 41 \u2014 Access to transport",
      "Buses, trains, stations, and airports must have ramps, accessible signs, and trained staff to help.",
      "A bus, train, or station did not let you board, had no ramp or lift, or had no accessible toilet.",
      "bus"),
    S("VIII", "Section 42 \u2014 Access to information and communication technology",
      "Websites, apps, TV, and phones must work for people with disabilities. Books, forms, and notices must come in accessible formats too.",
      "A government website, app, TV news, or printed form is not accessible to you.",
      "computer", ts=20),
    S("VIII", "Section 43 \u2014 Consumer goods",
      "Things sold in shops — like phones, TVs, and household goods — should be made so that persons with disabilities can use them.",
      "A product you bought is not usable because of your disability and the maker did not provide an accessible version.",
      "computer"),
    S("VIII", "Section 44 \u2014 Mandatory observance of accessibility norms",
      "No new building, transport, or service can get a permit unless it follows the accessibility rules.",
      "A new building or service has been approved but does not meet the accessibility standards.",
      "building_ramp"),
    S("VIII", "Section 45 \u2014 Time limit for making existing infrastructure and premises accessible and action for that purpose",
      "Old buildings and services must be made accessible within five years. This is the law, not a request.",
      "A public building or service is still inaccessible long after the law required it to be fixed.",
      "clock", ts=14),
    S("VIII", "Section 46 \u2014 Time limit for accessibility by service providers",
      "Service providers — like banks, hospitals, and shops — must make their services accessible within a fixed time set by the government.",
      "A service provider has not made its service accessible within the time the government set.",
      "clock", ts=20),
    S("VIII", "Section 47 \u2014 Human resource development",
      "The government must train enough teachers, doctors, and other staff to support persons with disabilities.",
      "There are no trained special educators, sign language interpreters, or rehab workers in your area.",
      "school"),
    S("VIII", "Section 48 \u2014 Social audit",
      "The government must check, at least every five years, if the schemes for persons with disabilities are working — and persons with disabilities themselves should be part of that check.",
      "A scheme is not being audited, or persons with disabilities were not part of the audit.",
      "clipboard_check"),

    # ---------- Chapter IX ----------
    S("IX", "Section 49 \u2014 Competent authority",
      "The government must name the office that handles registration of organisations working with persons with disabilities.",
      "(For reference. This sets up the registration system for disability organisations.)",
      "document"),
    S("IX", "Section 50 \u2014 Registration",
      "Any institution that wants to work with persons with disabilities must be registered. Without registration, they cannot legally provide services.",
      "An institution is offering disability services without being registered.",
      "document"),
    S("IX", "Section 51 \u2014 Application and grant of certificate of registration",
      "An institution must apply for a registration certificate. The competent authority gives the certificate after checking that the rules are met.",
      "Your application for a registration certificate is taking too long, or was unfairly refused.",
      "document", ts=20),
    S("IX", "Section 52 \u2014 Revocation of registration",
      "If a registered institution breaks the rules or harms persons with disabilities, its registration can be cancelled.",
      "A registered institution is treating persons with disabilities badly. You can ask the authority to cancel its registration.",
      "warning"),
    S("IX", "Section 53 \u2014 Appeal",
      "If your registration is refused or cancelled, you have the right to appeal to a higher authority.",
      "Your institution's registration was refused or cancelled and you want to challenge that decision.",
      "justice_scales"),
    S("IX", "Section 54 \u2014 Act not to apply to institutions established or maintained by Central or State Government",
      "Government-run institutions for persons with disabilities do not need to register under this chapter — different rules apply to them.",
      "(For reference. Government institutions are exempt from this registration requirement.)",
      "document", ts=14),
    S("IX", "Section 55 \u2014 Assistance to registered institutions",
      "The government can give grants and other help to registered organisations that serve persons with disabilities.",
      "(For reference. This section is the basis for funding disability organisations.)",
      "money_jar"),

    # ---------- Chapter X ----------
    S("X", "Section 56 \u2014 Guidelines for assessment of specified disabilities",
      "The government must publish clear guidelines on how to assess each of the 21 specified disabilities.",
      "(For reference. The assessment guidelines are issued by the central government.)",
      "clipboard_check"),
    S("X", "Section 57 \u2014 Designation of certifying authorities",
      "The government must name the doctors and offices that can issue a disability certificate.",
      "(For reference. This sets up the disability certification system.)",
      "document"),
    S("X", "Section 58 \u2014 Procedure for certification",
      "There must be a clear, fair process to get a disability certificate. The certificate is valid all over India.",
      "A doctor or office is making the certificate process hard, slow, or unclear.",
      "clipboard_check"),
    S("X", "Section 59 \u2014 Appeal against a decision of certifying authority",
      "If you are refused a disability certificate or get the wrong one, you can appeal.",
      "Your application for a disability certificate was refused, or the percentage was wrong.",
      "justice_scales", ts=20),

    # ---------- Chapter XI ----------
    S("XI", "Section 60 \u2014 Constitution of Central Advisory Board on Disability",
      "The central government must set up a Central Advisory Board on Disability to advise on disability matters.",
      "(For reference. This sets up the national disability advisory body.)",
      "committee_people", ts=22),
    S("XI", "Section 61 \u2014 Terms and conditions of service of members",
      "The rules say how Central Advisory Board members are paid and how long they serve.",
      "(For reference. Administrative rule for Central Advisory Board members.)",
      "document"),
    S("XI", "Section 62 \u2014 Disqualifications",
      "Some people cannot be members of the Central Advisory Board — for example, people convicted of serious crimes.",
      "(For reference. Lists who cannot be a Central Advisory Board member.)",
      "document"),
    S("XI", "Section 63 \u2014 Vacation of seats by Members",
      "If a Central Advisory Board member misses too many meetings or breaks the rules, they lose their seat.",
      "(For reference. Administrative rule for Central Advisory Board.)",
      "document"),
    S("XI", "Section 64 \u2014 Meetings of the Central Advisory Board on disability",
      "The Central Advisory Board must meet at least once every six months.",
      "(For reference. The Central Advisory Board has not met in over six months.)",
      "committee_people", ts=22),
    S("XI", "Section 65 \u2014 Functions of Central Advisory Board on disability",
      "The Central Advisory Board advises the central government on disability rules, schemes, and reports.",
      "(For reference. Lists what the Central Advisory Board does.)",
      "document", ts=22),
    S("XI", "Section 66 \u2014 State Advisory Board on disability",
      "Every state must set up its own State Advisory Board on Disability.",
      "Your state does not have a State Advisory Board on Disability.",
      "committee_people"),
    S("XI", "Section 67 \u2014 Terms and conditions of service of Members",
      "Rules for State Advisory Board members — pay, term, and conditions.",
      "(For reference. Administrative rule for State Advisory Board members.)",
      "document"),
    S("XI", "Section 68 \u2014 Disqualification",
      "Same kind of rules as the Central Board: who cannot be a State Advisory Board member.",
      "(For reference. Administrative rule for State Advisory Board.)",
      "document"),
    S("XI", "Section 69 \u2014 Vacation of seats",
      "When a State Advisory Board member loses their seat — missed meetings, broken rules, etc.",
      "(For reference. Administrative rule for State Advisory Board.)",
      "document"),
    S("XI", "Section 70 \u2014 Meetings of State Advisory Board on disability",
      "The State Advisory Board must meet at least once every six months.",
      "Your State Advisory Board has not met in over six months.",
      "committee_people", ts=22),
    S("XI", "Section 71 \u2014 Functions of State Advisory Board on disability",
      "The State Advisory Board advises the state government on disability rules, schemes, and reports.",
      "(For reference. Lists what the State Advisory Board does.)",
      "document", ts=22),
    S("XI", "Section 72 \u2014 District-level Committee on disability",
      "Every district must have a committee that handles disability matters at the district level.",
      "Your district has no disability committee, or the committee is not working.",
      "committee_people"),
    S("XI", "Section 73 \u2014 Vacancies not to invalidate proceedings",
      "Even if some seats on the Boards or committees are empty, their decisions are still valid.",
      "(For reference. Administrative rule.)",
      "document"),

    # ---------- Chapter XII ----------
    S("XII", "Section 74 \u2014 Appointment of Chief Commissioner and Commissioners",
      "The central government appoints a Chief Commissioner for Persons with Disabilities and two Commissioners under the Chief Commissioner.",
      "(For reference. This sets up the national disability complaints office.)",
      "commissioner_badge", ts=20),
    S("XII", "Section 75 \u2014 Functions of Chief Commissioner",
      "The Chief Commissioner watches over how the Act is being followed across India and takes action when rights are broken.",
      "(For reference. The Chief Commissioner is who you escalate to after the State Commissioner.)",
      "commissioner_badge"),
    S("XII", "Section 76 \u2014 Action of appropriate authorities on recommendation of Chief Commissioner",
      "When the Chief Commissioner makes a recommendation, the government office must act on it within a fixed time and report back.",
      "The Chief Commissioner gave a recommendation and the office did not act on it.",
      "document", ts=16),
    S("XII", "Section 77 \u2014 Powers of Chief Commissioner",
      "The Chief Commissioner can call witnesses, ask for documents, and act like a civil court when looking into a complaint.",
      "You filed a complaint with the Chief Commissioner and they have not used their powers.",
      "justice_scales"),
    S("XII", "Section 78 \u2014 Annual and special reports by Chief Commissioner",
      "The Chief Commissioner must publish a yearly report. The report goes to Parliament.",
      "The Chief Commissioner has not published a yearly report.",
      "document", ts=22),
    S("XII", "Section 79 \u2014 Appointment of State Commissioner in States",
      "Every state must appoint a State Commissioner for Persons with Disabilities.",
      "Your state does not have a State Commissioner for Persons with Disabilities.",
      "commissioner_badge"),
    S("XII", "Section 80 \u2014 Functions of State Commissioner",
      "The State Commissioner handles complaints in the state and watches over the Act in the state.",
      "Use the State Commissioner as your first stop for complaints inside your state.",
      "commissioner_badge"),
    S("XII", "Section 81 \u2014 Action by appropriate authorities on recommendation of State Commissioner",
      "When the State Commissioner makes a recommendation, the office must act on it and report back within a fixed time.",
      "The State Commissioner gave a recommendation and the office did not act on it.",
      "document", ts=16),
    S("XII", "Section 82 \u2014 Powers of State Commissioner",
      "The State Commissioner can call witnesses, ask for documents, and act like a civil court when looking into a complaint.",
      "You filed a complaint and the State Commissioner has not used their powers.",
      "justice_scales"),
    S("XII", "Section 83 \u2014 Annual and special reports by State Commissioner",
      "The State Commissioner must publish a yearly report. The report goes to the state government.",
      "The State Commissioner in your state has not published a yearly report.",
      "document", ts=22),

    # ---------- Chapter XIII ----------
    S("XIII", "Section 84 \u2014 Special Court",
      "Each state must have a Special Court to handle serious offences against persons with disabilities under this Act.",
      "A serious case under this Act is not being heard in a Special Court.",
      "courthouse"),
    S("XIII", "Section 85 \u2014 Special Public Prosecutor",
      "Each Special Court must have a Special Public Prosecutor to argue the case for the victim.",
      "A case is being heard in the Special Court without a Special Public Prosecutor.",
      "courthouse"),

    # ---------- Chapter XIV ----------
    S("XIV", "Section 86 \u2014 National Fund for persons with disabilities",
      "There is a National Fund for Persons with Disabilities. The Fund pays for schemes, support, and grants.",
      "(For reference. The Fund pays for many disability programmes.)",
      "money_jar"),
    S("XIV", "Section 87 \u2014 Accounts and audit",
      "The National Fund must keep clear accounts. The accounts must be audited every year.",
      "(For reference. Audit rule for the National Fund.)",
      "clipboard_check"),

    # ---------- Chapter XV ----------
    S("XV", "Section 88 \u2014 State Fund for persons with disabilities",
      "Every state must set up a State Fund for Persons with Disabilities, like the National Fund.",
      "Your state has no State Fund, or the State Fund is not being used for persons with disabilities.",
      "money_jar"),

    # ---------- Chapter XVI ----------
    S("XVI", "Section 89 \u2014 Punishment for contravention of provisions of Act or rules or regulations made thereunder",
      "If anyone breaks the rules of this Act, they can be fined. The first time: up to Rs 10,000. The next time: Rs 50,000 to Rs 5 lakh.",
      "Someone or some office is breaking the rules of this Act and you want them to be punished.",
      "warning", ts=14),
    S("XVI", "Section 90 \u2014 Offences by companies",
      "If a company breaks the rules of this Act, the people running the company can also be punished.",
      "A company has broken the rules and you want the managers to be held responsible too.",
      "warning"),
    S("XVI", "Section 91 \u2014 Punishment for fraudulently availing any benefit meant for persons with benchmark disabilities",
      "If a person fakes a disability to get benefits meant for persons with benchmark disabilities, they can be jailed for up to 2 years and fined up to Rs 1 lakh.",
      "Someone is faking a disability certificate to get a job, seat, or benefit reserved for persons with disabilities.",
      "warning", ts=14),
    S("XVI", "Section 92 \u2014 Punishment for offences of atrocities",
      "If someone insults you, hurts you, or humiliates you because of your disability, they can go to jail for 6 months to 5 years.",
      "Someone publicly insults, beats, denies food, or sexually abuses a person with a disability because of their disability.",
      "warning"),
    S("XVI", "Section 93 \u2014 Punishment for failure to furnish information",
      "If anyone refuses to give information that the Chief Commissioner or State Commissioner asks for, they can be fined.",
      "An office refused to send information to the Chief Commissioner or State Commissioner.",
      "warning"),
    S("XVI", "Section 94 \u2014 Previous sanction of appropriate Government",
      "Some cases under this Act need permission from the government before they can be filed in court.",
      "(For reference. Procedural rule for filing certain cases.)",
      "document"),
    S("XVI", "Section 95 \u2014 Alternative punishments",
      "The court can give other punishments — like community service or training — instead of fine or jail in some cases.",
      "(For reference. Court has flexibility on punishment in some cases.)",
      "justice_scales"),

    # ---------- Chapter XVII ----------
    S("XVII", "Section 96 \u2014 Application of other laws not barred",
      "This Act does not stop you from using other laws too. You can use this Act and other laws together.",
      "(For reference. You can use multiple legal protections at once.)",
      "document"),
    S("XVII", "Section 97 \u2014 Protection of action taken in good faith",
      "Government officers and others cannot be sued for actions they take honestly under this Act.",
      "(For reference. Protects officers acting in good faith.)",
      "document"),
    S("XVII", "Section 98 \u2014 Power to remove difficulties",
      "If something in this Act is unclear or hard to do, the central government can issue an order to fix it within three years of the Act starting.",
      "(For reference. Allows the central government to issue clarifications.)",
      "document"),
    S("XVII", "Section 99 \u2014 Power to amend Schedule",
      "The central government can change the list of 21 specified disabilities — add, remove, or change them.",
      "(For reference. The list of disabilities can be updated by the central government.)",
      "book_simple"),
    S("XVII", "Section 100 \u2014 Power of Central Government to make rules",
      "The central government can make rules to put this Act into action.",
      "(For reference. Many parts of the Act work through rules made under Section 100.)",
      "document"),
    S("XVII", "Section 101 \u2014 Power of State Government to make rules",
      "Each state government can make its own rules to put this Act into action in the state.",
      "(For reference. State governments make their own rules under Section 101.)",
      "document"),
    S("XVII", "Section 102 \u2014 Repeal and savings",
      "This Act replaces the older 1995 Persons with Disabilities Act, but actions already taken under the old Act are still valid.",
      "(For reference. The 1995 PWD Act is repealed and replaced by this 2016 Act.)",
      "book_simple"),
]


# ============================================================
# orchestration
# ============================================================

def build():
    os.makedirs(PROJECT, exist_ok=True)

    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    queue = []
    queue.append(("cover", None))
    queue.append(("about", None))
    queue.append(("how_to_use", None))
    queue.append(("map", None))

    seen = set()
    for s in SECTIONS:
        ch = s["chapter"]
        if ch not in seen:
            queue.append(("divider", CHAPTER_DIVIDERS[ch]))
            seen.add(ch)
        queue.append(("section", s))

    queue.append(("complain", None))
    queue.append(("help", None))
    queue.append(("summary", None))
    queue.append(("back", None))

    total = len(queue)

    for i, (kind, payload) in enumerate(queue, start=1):
        if kind == "cover":
            s_cover(prs, i, total)
        elif kind == "about":
            s_about(prs, i, total)
        elif kind == "how_to_use":
            s_how_to_use(prs, i, total)
        elif kind == "map":
            s_map(prs, i, total)
        elif kind == "divider":
            s_chapter_divider(prs, i, total, **payload)
        elif kind == "section":
            s_section(prs, i, total, **payload)
        elif kind == "complain":
            s_complain(prs, i, total)
        elif kind == "help":
            s_help(prs, i, total)
        elif kind == "summary":
            s_summary(prs, i, total)
        elif kind == "back":
            s_back(prs, i, total)

    try:
        prs.save(OUTPUT)
    except PermissionError:
        print(f"PermissionError: {OUTPUT} is open in PowerPoint. Close it and rerun.")
        raise
    print(f"OK {total} slides -> {OUTPUT}")


if __name__ == "__main__":
    build()
